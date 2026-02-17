#!/usr/bin/env python3
"""
Index TFK18 data into self-contained volumes directory.
Creates: TFK18/volumes/esdata, chroma, manifest
"""
import os
import sys
import json
import hashlib
import time
import requests
from typing import List, Tuple

# Paths relative to TFK18 root
TFK18_ROOT = "/media/felix/RAG/TFK18"
TFK18_VOLUMES = os.path.join(TFK18_ROOT, "volumes")
TFK18_DATA = TFK18_ROOT  # Data files directly in TFK18/

ES_URL = os.getenv("ES_URL", "http://localhost:9200")
ES_INDEX = "rag_tfk18_v1"  # Must match tenants/tfk18.yaml es_index
CHROMA_PATH = os.path.join(TFK18_VOLUMES, "chroma")
MANIFEST_PATH = os.path.join(TFK18_VOLUMES, "manifest", "manifest.sqlite3")

# Ensure directories exist
os.makedirs(os.path.join(TFK18_VOLUMES, "esdata"), exist_ok=True)
os.makedirs(CHROMA_PATH, exist_ok=True)
os.makedirs(os.path.join(TFK18_VOLUMES, "manifest"), exist_ok=True)
os.makedirs(os.path.join(TFK18_VOLUMES, "logs"), exist_ok=True)

EMBED_MODEL = os.getenv("EMBED_MODEL", "all-MiniLM-L6-v2")
ES_BATCH = int(os.getenv("ES_BATCH", "100"))
EMBED_BATCH = int(os.getenv("EMBED_BATCH", "128"))
CHROMA_BATCH = int(os.getenv("CHROMA_BATCH", "256"))

def chunk_text(text, size=1200, overlap=180):
    text = (text or "").strip()
    if not text:
        return []
    out = []
    i = 0
    n = len(text)
    step = max(1, size - overlap)
    while i < n:
        out.append(text[i:i+size])
        i += step
    return out

def stable_id(path, chunk_index):
    h = hashlib.sha1(path.encode("utf-8", errors="ignore")).hexdigest()[:16]
    return f"tfk18:{h}:{chunk_index}"

def get_file_hash(filepath):
    h = hashlib.sha256()
    try:
        with open(filepath, "rb") as f:
            while chunk := f.read(8192):
                h.update(chunk)
        return h.hexdigest()
    except:
        return None

def scan_files(root_dir):
    """Recursively scan all files in directory"""
    files = []
    for dirpath, dirnames, filenames in os.walk(root_dir):
        for filename in filenames:
            filepath = os.path.join(dirpath, filename)
            files.append(filepath)
    return files

def extract_text_basic(filepath):
    """Basic text extraction for various file types"""
    ext = os.path.splitext(filepath)[1].lower()
    
    try:
        if ext in ['.txt', '.md', '.py', '.js', '.html', '.xml', '.json', '.csv']:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext in ['.pdf']:
            # Use pdfplumber if available
            try:
                import pdfplumber
                with pdfplumber.open(filepath) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            except:
                return None
        elif ext in ['.docx']:
            try:
                from docx import Document
                doc = Document(filepath)
                return "\n".join(p.text for p in doc.paragraphs)
            except:
                return None
        elif ext in ['.pptx']:
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
            except:
                return None
        elif ext in ['.eml', '.msg']:
            # Use existing text_loaders
            sys.path.insert(0, '/media/felix/RAG/AGENTIC/indexer/app')
            from text_loaders import read_eml_with_attachments, read_msg
            if ext == '.eml':
                return read_eml_with_attachments(filepath)
            else:
                return read_msg(filepath)
        else:
            return None
    except Exception as e:
        print(f"⚠️ Error extracting {filepath}: {e}")
        return None

def es_create_index():
    """Create ES index if not exists"""
    url = f"{ES_URL}/{ES_INDEX}"
    try:
        r = requests.head(url, timeout=10)
        if r.status_code == 200:
            print(f"✅ ES index {ES_INDEX} already exists")
            return
    except:
        pass
    
    # Create index with mapping matching rag_files_v1 structure (nested fields)
    mapping = {
        "mappings": {
            "properties": {
                "content": {"type": "text"},
                "path": {
                    "properties": {
                        "virtual": {
                            "type": "keyword",
                            "fields": {
                                "fulltext": {"type": "text"}
                            }
                        }
                    }
                },
                "file": {
                    "properties": {
                        "filename": {"type": "keyword", "store": True},
                        "extension": {"type": "keyword"},
                        "content_type": {"type": "keyword"},
                        "filesize": {"type": "long"},
                        "checksum": {"type": "keyword"}
                    }
                },
                "indexed_at": {"type": "date"}
            }
        }
    }
    
    try:
        r = requests.put(url, json=mapping, timeout=30)
        if r.status_code in [200, 201]:
            print(f"✅ Created ES index: {ES_INDEX}")
        else:
            print(f"⚠️ ES index creation: {r.status_code} - {r.text}")
    except Exception as e:
        print(f"❌ Failed to create ES index: {e}")

def es_index_batch(docs):
    """Bulk index to ES"""
    if not docs:
        return 0
    
    url = f"{ES_URL}/{ES_INDEX}/_bulk"
    lines = []
    for doc in docs:
        lines.append(json.dumps({"index": {"_index": ES_INDEX}}))
        lines.append(json.dumps(doc))
    
    body = "\n".join(lines) + "\n"
    
    try:
        headers = {"Content-Type": "application/json"}
        r = requests.post(url, data=body, headers=headers, timeout=60)
        if r.status_code == 200:
            result = r.json()
            errors = result.get("errors", False)
            if errors:
                print(f"⚠️ ES bulk had errors: {result.get('items', [])[:2]}")
            return len(docs)
        else:
            print(f"⚠️ ES bulk failed: {r.status_code}")
            return 0
    except Exception as e:
        print(f"⚠️ ES bulk error: {e}")
        return 0

def batch_encode(embedder, chunks: List[str]) -> List[List[float]]:
    if not chunks:
        return []
    embeddings = embedder.encode(chunks, show_progress_bar=False, convert_to_numpy=True)
    return embeddings.tolist()

INDEXABLE_EXTS = {'.pdf', '.docx', '.doc', '.eml', '.msg', '.txt', '.md', '.csv',
                  '.xlsx', '.xls', '.pptx', '.ppt', '.html', '.xml', '.json'}

def main():
    
    print("=" * 60)
    print("TFK18 → Elasticsearch Index Builder")
    print("=" * 60)
    print(f"Data: {TFK18_DATA}")
    print(f"ES Index: {ES_INDEX}")
    print(f"ES URL: {ES_URL}")
    print("=" * 60)
    
    # Scan files (only indexable types)
    print("\n🔍 Scanning files...")
    all_files = scan_files(TFK18_DATA)
    all_files = [f for f in all_files if not f.startswith(TFK18_VOLUMES)]
    files = [f for f in all_files if os.path.splitext(f)[1].lower() in INDEXABLE_EXTS]
    print(f"📁 Found {len(all_files):,} total files, {len(files):,} indexable")
    
    if not files:
        print("❌ No indexable files found!")
        return
    
    # Setup ES
    es_create_index()
    
    # Process files (ES only – Chroma can be added later)
    print(f"\n🚀 Processing {len(files):,} files...")
    es_docs = []
    processed = 0
    skipped = 0
    es_indexed = 0
    start_time = time.time()
    
    for filepath in files:
        rel_path = os.path.relpath(filepath, TFK18_DATA)
        filename = os.path.basename(filepath)
        ext = os.path.splitext(filename)[1].lower()
        
        # Extract text
        content = extract_text_basic(filepath)
        if not content:
            continue
        
        file_hash = get_file_hash(filepath)
        
        # ES document (nested structure matching rag_files_v1)
        es_doc = {
            "content": content,
            "file": {
                "filename": filename,
                "extension": ext.lstrip("."),
                "filesize": os.path.getsize(filepath),
                "checksum": file_hash,
            },
            "path": {
                "virtual": "/" + rel_path,
            },
        }
        es_docs.append(es_doc)
        
        processed += 1
        if processed % 200 == 0:
            elapsed = time.time() - start_time
            rate = processed / elapsed if elapsed > 0 else 0
            print(f"  ⏳ {processed:,}/{len(files):,} files ({rate:.0f}/s), ES: {es_indexed:,} docs")
        
        # Batch ES index
        if len(es_docs) >= ES_BATCH:
            indexed = es_index_batch(es_docs)
            es_indexed += indexed
            es_docs = []
    
    # Final ES batch
    if es_docs:
        indexed = es_index_batch(es_docs)
        es_indexed += indexed
    
    elapsed = time.time() - start_time
    
    print("\n" + "=" * 60)
    print("✅ TFK18 ES Index Complete:")
    print(f"   Files scanned: {len(files):,}")
    print(f"   Files with text: {processed:,} (skipped {len(files) - processed:,})")
    print(f"   ES docs indexed: {es_indexed:,}")
    print(f"   ES Index: {ES_INDEX}")
    print(f"   Duration: {elapsed:.0f}s ({elapsed/60:.1f}min)")
    print("=" * 60)

if __name__ == "__main__":
    main()
