import os
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import extract_msg


def read_text_file(path: str) -> str:
    """
    Robust text reader with encoding fallbacks.
    """
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            with open(path, "r", encoding=enc, errors="strict") as f:
                return f.read()
        except Exception:
            pass
    with open(path, "rb") as f:
        return f.read().decode("utf-8", errors="ignore")


def read_text_bytes(b: bytes) -> str:
    """
    Decode bytes with robust fallbacks.
    """
    for enc in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
        try:
            return b.decode(enc)
        except Exception:
            pass
    return b.decode("utf-8", errors="ignore")


def read_docx(path: str) -> str:
    """
    Extract text from DOCX paragraphs.
    """
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text]
    return "\n".join(parts)


def read_pptx(path: str) -> str:
    """
    Extract text from PPTX slides/shapes.
    """
    prs = Presentation(path)
    out = []
    for si, slide in enumerate(prs.slides, 1):
        out.append(f"--- Slide {si} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                out.append(shape.text)
    return "\n".join(out)


def read_xlsx(path: str) -> str:
    """
    Extract all sheets to CSV-like text.
    """
    xl = pd.ExcelFile(path, engine="openpyxl")
    out = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        out.append(f"--- Sheet: {sheet} ---")
        out.append(df.to_csv(index=False))
    return "\n".join(out)


def read_html(path: str) -> str:
    """
    Extract visible text from HTML.
    """
    raw = read_text_file(path)
    soup = BeautifulSoup(raw, "lxml")
    return soup.get_text("\n")


def read_msg(path: str) -> str:
    """
    Extract subject/from/to/date/body from .msg (Outlook).
    """
    m = extract_msg.Message(path)
    m.process()

    fields = []
    if getattr(m, "subject", None):
        fields.append(f"Subject: {m.subject}")
    if getattr(m, "sender", None):
        fields.append(f"From: {m.sender}")
    if getattr(m, "to", None):
        fields.append(f"To: {m.to}")
    if getattr(m, "date", None):
        fields.append(f"Date: {m.date}")

    body = getattr(m, "body", "") or ""
    return "\n".join(fields + ["", body]).strip()


def extract_pdf_metadata(path: str) -> dict:
    """Extract metadata from PDF (author, title, creation date, producer)."""
    meta = {}
    try:
        import fitz
        doc = fitz.open(path)
        info = doc.metadata or {}
        if info.get("author"):
            meta["author"] = info["author"][:200]
        if info.get("title"):
            meta["title"] = info["title"][:300]
        if info.get("creationDate"):
            meta["creation_date"] = info["creationDate"][:50]
        if info.get("producer"):
            meta["producer"] = info["producer"][:200]
        if info.get("subject"):
            meta["subject"] = info["subject"][:300]
        meta["page_count"] = len(doc)
        doc.close()
    except Exception as e:
        meta["_meta_error"] = str(e)[:100]
    return meta


def extract_docx_metadata(path: str) -> dict:
    """Extract metadata from DOCX (author, title, created, modified)."""
    meta = {}
    try:
        doc = Document(path)
        props = doc.core_properties
        if props.author:
            meta["author"] = str(props.author)[:200]
        if props.title:
            meta["title"] = str(props.title)[:300]
        if props.subject:
            meta["subject"] = str(props.subject)[:300]
        if props.created:
            meta["creation_date"] = props.created.isoformat()
        if props.modified:
            meta["modified_date"] = props.modified.isoformat()
        if props.last_modified_by:
            meta["last_modified_by"] = str(props.last_modified_by)[:200]
    except Exception as e:
        meta["_meta_error"] = str(e)[:100]
    return meta


def extract_eml_metadata(path: str) -> dict:
    """Extract email headers as structured metadata from EML file."""
    import email
    import email.policy
    meta = {}
    try:
        with open(path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=email.policy.default)
        if msg.get("Subject"):
            meta["email_subject"] = str(msg["Subject"])[:300]
        if msg.get("From"):
            meta["email_from"] = str(msg["From"])[:200]
        if msg.get("To"):
            meta["email_to"] = str(msg["To"])[:500]
        if msg.get("Date"):
            meta["email_date"] = str(msg["Date"])[:50]
        if msg.get("Cc"):
            meta["email_cc"] = str(msg["Cc"])[:500]
    except Exception as e:
        meta["_meta_error"] = str(e)[:100]
    return meta


def extract_msg_metadata(path: str) -> dict:
    """Extract email headers as structured metadata from MSG file."""
    meta = {}
    try:
        m = extract_msg.Message(path)
        m.process()
        if getattr(m, "subject", None):
            meta["email_subject"] = str(m.subject)[:300]
        if getattr(m, "sender", None):
            meta["email_from"] = str(m.sender)[:200]
        if getattr(m, "to", None):
            meta["email_to"] = str(m.to)[:500]
        if getattr(m, "date", None):
            meta["email_date"] = str(m.date)[:50]
        if getattr(m, "cc", None):
            meta["email_cc"] = str(m.cc)[:500]
    except Exception as e:
        meta["_meta_error"] = str(e)[:100]
    return meta


def read_eml_with_attachments(path: str) -> dict:
    """
    Extract text from .eml file including attachments with OCR.
    Returns: {"text": "combined text", "attachments": [{"filename": ..., "text": ...}]}
    """
    import email
    import email.policy
    from email import message_from_binary_file
    from io import BytesIO
    
    # Try to import OCR libraries
    try:
        import pytesseract
        from PIL import Image
        OCR_AVAILABLE = True
    except ImportError:
        OCR_AVAILABLE = False
    
    # Try to import PDF reader
    try:
        import fitz  # PyMuPDF
        PDF_AVAILABLE = True
    except ImportError:
        PDF_AVAILABLE = False

    def extract_text_from_attachment(part):
        """Extract text from attachment with OCR if needed"""
        filename = part.get_filename() or "unknown"
        content_type = part.get_content_type()
        payload = part.get_payload(decode=True) or b""
        
        if not payload:
            return {"filename": filename, "text": "", "content_type": content_type}
        
        text = ""
        ext = os.path.splitext(filename)[1].lower()
        
        # Text files
        if content_type.startswith("text/") or ext in [".txt", ".csv", ".json", ".xml", ".html"]:
            text = read_text_bytes(payload)
        
        # PDF files
        elif ext == ".pdf" and PDF_AVAILABLE:
            try:
                doc = fitz.open(stream=payload, filetype="pdf")
                for page in doc:
                    text += page.get_text()
            except Exception as e:
                text = f"[PDF extraction error: {e}]"
        
        # Images with OCR
        elif ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff"] and OCR_AVAILABLE:
            try:
                img = Image.open(BytesIO(payload))
                text = pytesseract.image_to_string(img)
            except Exception as e:
                text = f"[OCR error: {e}]"
        
        # DOCX files
        elif ext == ".docx":
            try:
                from docx import Document
                doc = Document(BytesIO(payload))
                text = "\n".join([p.text for p in doc.paragraphs])
            except Exception as e:
                text = f"[DOCX extraction error: {e}]"
        
        return {"filename": filename, "text": text, "content_type": content_type}

    # Parse EML file
    with open(path, "rb") as f:
        msg = message_from_binary_file(f, policy=email.policy.default)
    
    # Extract headers
    headers = []
    for header in ["Subject", "From", "To", "Date", "Cc", "Bcc"]:
        value = msg.get(header)
        if value:
            headers.append(f"{header}: {value}")
    
    # Extract body and attachments
    body_text = ""
    attachments = []
    
    if msg.is_multipart():
        for part in msg.walk():
            content_type = part.get_content_type()
            
            # Skip container parts
            if content_type.startswith("multipart/"):
                continue
            
            # Extract body from text parts
            if content_type.startswith("text/") and not part.get_filename():
                try:
                    body_text += part.get_content()
                except:
                    pass
            
            # Extract attachments
            elif part.get_filename():
                att = extract_text_from_attachment(part)
                attachments.append(att)
    else:
        # Single part message
        try:
            body_text = msg.get_content()
        except:
            pass
    
    # Combine all text
    attachment_text = "\n\n".join([
        f"--- Attachment: {a['filename']} ---\n{a['text']}"
        for a in attachments if a.get("text")
    ])
    
    full_text = "\n\n".join([
        "\n".join(headers),
        "",
        body_text,
        attachment_text
    ]).strip()
    
    return {
        "text": full_text,
        "attachments": attachments
    }
