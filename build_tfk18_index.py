#!/usr/bin/env python3
"""
TFK18 Self-Contained Index Builder
Creates isolated index under /media/felix/RAG/TFK18/volumes/
Uses separate ES instance on port 9201 - does NOT touch main index on 9200
"""
import os
import sys
import time
import hashlib
import subprocess
import signal
import requests
from pathlib import Path

# Configuration
TFK18_ROOT = "/media/felix/RAG/TFK18"
TFK18_VOLUMES = os.path.join(TFK18_ROOT, "volumes")
TFK18_DATA = TFK18_ROOT

ES_PORT = 9201  # Separate port - avoids conflict with main ES on 9200
ES_INDEX = "rag_files_v1"  # Can use same name since ES is isolated
CHROMA_PATH = os.path.join(TFK18_VOLUMES, "chroma")

# Ensure directories exist (but don't create if TFK18 doesn't exist)
if not os.path.exists(TFK18_ROOT):
    print(f"❌ TFK18 directory not found: {TFK18_ROOT}")
    print("Please create /media/felix/RAG/TFK18 with your data first")
    sys.exit(1)

os.makedirs(os.path.join(TFK18_VOLUMES, "esdata"), exist_ok=True)
os.makedirs(CHROMA_PATH, exist_ok=True)
os.makedirs(os.path.join(TFK18_VOLUMES, "logs"), exist_ok=True)

ES_URL = f"http://localhost:{ES_PORT}"
EMBED_MODEL = os.getenv("EMBED_MODEL", "all-MiniLM-L6-v2")
EMBED_BATCH = 128
CHROMA_BATCH = 256

# Colors for output
GREEN = "\033[92m"
YELLOW = "\033[93m"
RED = "\033[91m"
BLUE = "\033[94m"
RESET = "\033[0m"

def log(msg, color=RESET):
    print(f"{color}{msg}{RESET}")

def chunk_text(text, size=1200, overlap=180):
    text = (text or "").strip()
    if not text:
        return []
    out = []
    i = 0
    n = len(text)
    step = max(1, size - overlap)
    while i < n:
        out.append(text[i:i+size])
        i += step
    return out

def stable_id(path, chunk_index):
    h = hashlib.sha1(path.encode("utf-8", errors="ignore")).hexdigest()[:16]
    return f"tfk18:{h}:{chunk_index}"

def get_file_hash(filepath):
    h = hashlib.sha256()
    try:
        with open(filepath, "rb") as f:
            while chunk := f.read(8192):
                h.update(chunk)
        return h.hexdigest()
    except:
        return None

def scan_files(root_dir):
    """Recursively scan all files"""
    files = []
    for dirpath, dirnames, filenames in os.walk(root_dir):
        # Skip volumes directory
        if "volumes" in dirpath:
            continue
        for filename in filenames:
            filepath = os.path.join(dirpath, filename)
            files.append(filepath)
    return files

def extract_text_basic(filepath):
    """Extract text from various file types"""
    ext = os.path.splitext(filepath)[1].lower()
    
    try:
        if ext in ['.txt', '.md', '.py', '.js', '.html', '.xml', '.json', '.csv', '.log']:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == '.pdf':
            try:
                import pdfplumber
                with pdfplumber.open(filepath) as pdf:
                    return "\n".join(page.extract_text() or "" for page in pdf.pages)
            except:
                return None
        elif ext == '.docx':
            try:
                from docx import Document
                doc = Document(filepath)
                return "\n".join(p.text for p in doc.paragraphs)
            except:
                return None
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
            except:
                return None
        elif ext in ['.eml', '.msg']:
            sys.path.insert(0, '/media/felix/RAG/AGENTIC/indexer/app')
            from text_loaders import read_eml_with_attachments, read_msg
            if ext == '.eml':
                result = read_eml_with_attachments(filepath)
                return result.get('text', '') if isinstance(result, dict) else result
            else:
                return read_msg(filepath)
        elif ext == '.xlsx':
            try:
                from openpyxl import load_workbook
                wb = load_workbook(filepath, data_only=True)
                texts = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    sheet_text = f"--- Sheet: {sheet} ---\n"
                    for row in ws.iter_rows(values_only=True):
                        row_text = " | ".join(str(cell) for cell in row if cell)
                        if row_text.strip():
                            sheet_text += row_text + "\n"
                    texts.append(sheet_text)
                return "\n\n".join(texts)
            except Exception as e:
                return f"[XLSX metadata only - error: {e}]"
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
            # Images - try OCR if available, else metadata only
            try:
                import pytesseract
                from PIL import Image
                img = Image.open(filepath)
                text = pytesseract.image_to_string(img)
                if text.strip():
                    return text
                return "[Image - no extractable text]"
            except:
                return "[Image file - metadata only]"
        elif ext in ['.dwg', '.dxf', '.step', '.stp', '.iges', '.igs']:
            # CAD files - metadata only
            return f"[CAD file: {ext}]"
        elif ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
            # Archives - metadata only  
            return f"[Archive: {ext}]"
        else:
            # Unknown extension - return empty to trigger metadata-only indexing
            return None
    except Exception as e:
        log(f" Error extracting {filepath}: {e}", YELLOW)
        log(f"⚠️ Error extracting {filepath}: {e}", YELLOW)
        return None

def start_elasticsearch():
    """Start isolated ES instance for TFK18"""
    log("\n" + "="*60, BLUE)
    log("Starting isolated Elasticsearch for TFK18...", BLUE)
    log(f"Port: {ES_PORT}", BLUE)
    log(f"Data: {TFK18_VOLUMES}/esdata", BLUE)
    log("="*60 + "\n", BLUE)
    
    # Kill any existing ES on this port
    subprocess.run("pkill -f 'elasticsearch.*9201' 2>/dev/null || true", shell=True)
    time.sleep(2)
    
    # Start ES in Docker with isolated data
    es_container_name = "agentic-elasticsearch-tfk18"
    
    # Remove old container if exists
    subprocess.run(f"docker rm -f {es_container_name} 2>/dev/null || true", shell=True)
    
    # Start new ES container
    cmd = [
        "docker", "run", "-d",
        "--name", es_container_name,
        "-p", f"{ES_PORT}:9200",
        "-v", f"{TFK18_VOLUMES}/esdata:/usr/share/elasticsearch/data",
        "-e", "discovery.type=single-node",
        "-e", "xpack.security.enabled=false",
        "-e", "xpack.security.http.ssl.enabled=false",
        "-e", "ES_JAVA_OPTS=-Xms2g -Xmx2g",
        "--user", "1000:1000",
        "docker.elastic.co/elasticsearch/elasticsearch:8.12.2"
    ]
    
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        log(f"❌ Failed to start ES: {result.stderr}", RED)
        return False
    
    # Wait for ES to be ready
    log("⏳ Waiting for ES to be ready...")
    for i in range(30):
        try:
            r = requests.get(f"{ES_URL}/_cluster/health", timeout=5)
            if r.status_code == 200:
                log(f"✅ ES is ready! ({r.json().get('status', 'unknown')})", GREEN)
                return True
        except:
            pass
        time.sleep(2)
        log(f"  Attempt {i+1}/30...", YELLOW)
    
    log("❌ ES failed to start within timeout", RED)
    return False

def stop_elasticsearch():
    """Stop isolated ES instance"""
    log("\n🛑 Stopping isolated ES...", BLUE)
    subprocess.run("docker rm -f agentic-elasticsearch-tfk18 2>/dev/null || true", shell=True)
    log("✅ ES stopped", GREEN)

def es_create_index():
    """Create ES index"""
    url = f"{ES_URL}/{ES_INDEX}"
    try:
        r = requests.head(url, timeout=10)
        if r.status_code == 200:
            log(f"✅ ES index {ES_INDEX} already exists")
            return
    except:
        pass
    
    mapping = {
        "mappings": {
            "properties": {
                "content": {"type": "text", "analyzer": "standard"},
                "file.filename": {"type": "keyword"},
                "file.extension": {"type": "keyword"},
                "path.virtual": {"type": "keyword"},
                "file.hash": {"type": "keyword"},
                "indexed_at": {"type": "date"}
            }
        }
    }
    
    try:
        r = requests.put(url, json=mapping, timeout=30)
        if r.status_code in [200, 201]:
            log(f"✅ Created ES index: {ES_INDEX}", GREEN)
        else:
            log(f"⚠️ ES index creation: {r.status_code}", YELLOW)
    except Exception as e:
        log(f"❌ Failed to create ES index: {e}", RED)

def es_index_batch(docs):
    """Bulk index to ES"""
    if not docs:
        return 0
    
    url = f"{ES_URL}/{ES_INDEX}/_bulk"
    lines = []
    for doc in docs:
        lines.append('{"index":{}}')
        import json
        lines.append(json.dumps(doc))
    
    body = "\n".join(lines) + "\n"
    
    try:
        headers = {"Content-Type": "application/json"}
        r = requests.post(url, data=body, headers=headers, timeout=60)
        if r.status_code == 200:
            return len(docs)
        return 0
    except Exception as e:
        log(f"⚠️ ES bulk error: {e}", YELLOW)
        return 0

def batch_encode(embedder, chunks):
    if not chunks:
        return []
    embeddings = embedder.encode(chunks, show_progress_bar=False, convert_to_numpy=True)
    return embeddings.tolist()

def main():
    log("="*60, BLUE)
    log("TFK18 Self-Contained Index Builder", BLUE)
    log("="*60, BLUE)
    log(f"Data: {TFK18_DATA}")
    log(f"Volumes: {TFK18_VOLUMES}")
    log(f"ES: {ES_URL}")
    log(f"Index: {ES_INDEX}")
    log(f"Chroma: {CHROMA_PATH}")
    log("="*60, BLUE)
    
    # Verify isolation
    main_es_path = "/media/felix/RAG/1/volumes/esdata"
    if os.path.exists(main_es_path):
        log(f"\n✅ Main index protected at: {main_es_path}", GREEN)
        log(f"   TFK18 using separate path: {TFK18_VOLUMES}/esdata", GREEN)
    
    # Setup imports
    sys.path.insert(0, '/media/felix/RAG/AGENTIC/venv/lib/python3.12/site-packages')
    import chromadb
    from sentence_transformers import SentenceTransformer
    import json
    
    # Scan files
    log("\n🔍 Scanning TFK18 files...")
    files = scan_files(TFK18_DATA)
    log(f"📁 Found {len(files):,} files")
    
    if not files:
        log("❌ No files found!", RED)
        return
    
    # Start isolated ES
    if not start_elasticsearch():
        return
    
    try:
        # Setup Chroma (fresh)
        client = chromadb.PersistentClient(CHROMA_PATH)
        try:
            client.delete_collection("documents")
            log("🗑️ Deleted old Chroma collection")
        except:
            pass
        
        col = client.create_collection("documents")
        log(f"✅ Created Chroma collection", GREEN)
        
        # Create ES index
        es_create_index()
        
        # Load embedder
        log(f"\n🔄 Loading embedder: {EMBED_MODEL}")
        embedder = SentenceTransformer(EMBED_MODEL)
        
        # Process files
        log(f"\n🚀 Processing {len(files):,} files...")
        es_docs = []
        chroma_batches = []
        processed = 0
        es_indexed = 0
        chroma_chunks = 0
        
        for filepath in files:
            rel_path = os.path.relpath(filepath, TFK18_DATA)
            filename = os.path.basename(filepath)
            ext = os.path.splitext(filename)[1].lower()
            
            # Extract text (if possible)
            content = extract_text_basic(filepath)
            
            # Get file stats
            file_hash = get_file_hash(filepath)
            file_size = os.path.getsize(filepath)
            mtime = os.path.getmtime(filepath)
            
            # For ES: always index, even if only metadata
            es_doc = {
                "content": content or f"[File: {filename}]",
                "file.filename": filename,
                "file.extension": ext,
                "path.virtual": rel_path,
                "file.hash": file_hash,
                "file.size": file_size,
                "file.mtime": mtime
            }
            es_docs.append(es_doc)
            
            # For Chroma: only if we have actual content
            if content and len(content) > 50:
                chunks = chunk_text(content)
                for i, chunk in enumerate(chunks):
                    doc_id = stable_id(rel_path, i)
                    chroma_batches.append({
                        "id": doc_id,
                        "text": chunk,
                        "meta": {
                            "path": rel_path,
                            "filename": filename,
                            "chunk_index": i,
                            "extension": ext
                        }
                    })
            
            processed += 1
            
            # Batch ES index (every 100 docs)
            if len(es_docs) >= 100:
                indexed = es_index_batch(es_docs)
                es_indexed += indexed
                es_docs = []
            
            # Batch Chroma (every EMBED_BATCH chunks)
            if len(chroma_batches) >= EMBED_BATCH:
                texts = [b["text"] for b in chroma_batches[:EMBED_BATCH]]
                embeddings = batch_encode(embedder, texts)
                
                batch = chroma_batches[:EMBED_BATCH]
                try:
                    col.add(
                        ids=[b["id"] for b in batch],
                        documents=[b["text"] for b in batch],
                        metadatas=[b["meta"] for b in batch],
                        embeddings=embeddings
                    )
                    chroma_chunks += len(batch)
                except Exception as e:
                    log(f"⚠️ Chroma error: {e}", YELLOW)
                
                chroma_batches = chroma_batches[EMBED_BATCH:]
                
                if processed % 100 == 0:
                    log(f"  Progress: {processed:,} files | ES: {es_indexed:,} | Chroma: {chroma_chunks:,} chunks")
        
        # Final ES batch
        if es_docs:
            indexed = es_index_batch(es_docs)
            es_indexed += indexed
        
        # Final Chroma batches
        while chroma_batches:
            batch_size = min(EMBED_BATCH, len(chroma_batches))
            batch = chroma_batches[:batch_size]
            texts = [b["text"] for b in batch]
            embeddings = batch_encode(embedder, texts)
            
            try:
                col.add(
                    ids=[b["id"] for b in batch],
                    documents=[b["text"] for b in batch],
                    metadatas=[b["meta"] for b in batch],
                    embeddings=embeddings
                )
                chroma_chunks += len(batch)
            except Exception as e:
                log(f"⚠️ Chroma error: {e}", YELLOW)
            
            chroma_batches = chroma_batches[batch_size:]
        
        final_chroma_count = col.count()
        
        # Summary
        log("\n" + "="*60, GREEN)
        log("✅ TFK18 Index Complete!", GREEN)
        log("="*60, GREEN)
        log(f"Files processed: {processed:,}")
        log(f"ES docs: {es_indexed:,}")
        log(f"Chroma chunks: {chroma_chunks:,} (verified: {final_chroma_count:,})")
        log("\n📂 Locations:")
        log(f"   ES data: {TFK18_VOLUMES}/esdata/")
        log(f"   Chroma:  {CHROMA_PATH}/")
        log(f"   Logs:    {TFK18_VOLUMES}/logs/")
        log("\n🔒 Main index remained untouched at:")
        log(f"   /media/felix/RAG/1/volumes/")
        log("="*60, GREEN)
        
    finally:
        # Always stop ES
        stop_elasticsearch()
        
        log("\n📋 To switch to TFK18 later:")
        log("   docker-compose down")
        log("   mv /media/felix/RAG/1 /media/felix/RAG/1_backup")
        log("   mv /media/felix/RAG/TFK18 /media/felix/RAG/1")
        log("   docker-compose up -d")
        log("\n   (volumes/ moves with the folder - no config changes needed)")

if __name__ == "__main__":
    # Handle Ctrl+C gracefully
    def signal_handler(sig, frame):
        log("\n\n⚠️ Interrupted - cleaning up...", YELLOW)
        stop_elasticsearch()
        sys.exit(1)
    
    signal.signal(signal.SIGINT, signal_handler)
    signal.signal(signal.SIGTERM, signal_handler)
    
    main()
